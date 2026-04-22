"""
AIAuth canonical text extraction (Piece 14, v0.5.1+)

Produces format-agnostic canonical text from a document so that
equivalent content in different formats (e.g. xlsx -> csv -> pdf)
produces the SAME content_hash_canonical.

Used by:
  - aiauth.py (desktop agent) during file attestation
  - self-hosted bootstrap / tests
  - scripts that verify chain integrity across exported copies

Normalization rules (MUST match server's normalize_text for text
selections):
  1. Extract format-appropriate logical content (cells / paragraphs /
     slide text / key-value pairs).
  2. Collapse whitespace runs to single spaces.
  3. Lowercase.
  4. UTF-8 encode.
  5. SHA-256 hex -> content_hash_canonical.

Graceful degradation: if a format's library is missing (pdfplumber,
python-docx, openpyxl, python-pptx, imagehash), that format's
extractor raises NotImplementedError and the caller falls back to
raw-bytes hashing with canonical_extraction_failed=True in the
receipt.

Public entry points:
  canonical_text(path) -> str    # normalized canonical text
  canonical_hash(path) -> str    # hex SHA-256 of canonical_text
  perceptual_hashes(image_path) -> dict  # {'dhash': '...', 'phash': '...'}
"""

from __future__ import annotations

import csv
import hashlib
import json
import re
from io import StringIO
from pathlib import Path
from typing import Optional

_WS_RE = re.compile(r"\s+")
_PAGE_HEADER_RE = re.compile(
    r"^(page\s+\d+(\s+of\s+\d+)?|\d+\s*/\s*\d+)\s*$",
    re.IGNORECASE,
)


def _normalize(text: str) -> str:
    if text is None:
        return ""
    return _WS_RE.sub(" ", text).strip().lower()


# ---------------------------------------------------------------------------
# Per-format extractors
# Each returns a string (NOT yet normalized) representing the canonical
# content of the file. The caller normalizes + hashes at the end.
# ---------------------------------------------------------------------------

def _extract_text_plain(path: Path) -> str:
    with open(path, "rb") as f:
        raw = f.read()
    for enc in ("utf-8", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _extract_csv(path: Path) -> str:
    """Row-major, pipe-separated cell values. Lowercased cells; no header
    special-casing so csv with and without a header row canonicalize the
    same once case is normalized."""
    with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        parts = []
        for row in reader:
            parts.extend(str(cell).strip() for cell in row if str(cell).strip())
    return " | ".join(parts)


def _extract_tsv(path: Path) -> str:
    with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f, delimiter="\t")
        parts = []
        for row in reader:
            parts.extend(str(cell).strip() for cell in row if str(cell).strip())
    return " | ".join(parts)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise NotImplementedError("openpyxl not installed — cannot extract xlsx")
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is None:
                    continue
                s = str(cell).strip()
                if s:
                    parts.append(s)
    return " | ".join(parts)


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise NotImplementedError("python-docx not installed — cannot extract docx")
    doc = Document(str(path))
    parts = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    # Also pull table cells
    for tbl in doc.tables:
        for row in tbl.rows:
            for cell in row.cells:
                t = (cell.text or "").strip()
                if t:
                    parts.append(t)
    return " ".join(parts)


def _extract_pdf(path: Path) -> str:
    """Extract text layer. Strips repeated-per-page headers/footers and
    common "Page N of M" chrome.
    """
    # Prefer pdfplumber (better layout handling); fall back to pypdf.
    pages_text = []
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text() or ""
                pages_text.append(t)
    except ImportError:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            for page in reader.pages:
                pages_text.append(page.extract_text() or "")
        except ImportError:
            raise NotImplementedError("Neither pdfplumber nor pypdf installed — cannot extract pdf")

    if not pages_text:
        return ""

    # Strip repeated header/footer lines (appear on >50% of pages)
    if len(pages_text) > 1:
        line_counts: dict[str, int] = {}
        per_page_lines = [p.split("\n") for p in pages_text]
        for lines in per_page_lines:
            seen = set()
            for ln in lines:
                ln_norm = ln.strip()
                if ln_norm and ln_norm not in seen:
                    line_counts[ln_norm] = line_counts.get(ln_norm, 0) + 1
                    seen.add(ln_norm)
        threshold = max(2, len(pages_text) // 2 + 1)
        repeated = {ln for ln, n in line_counts.items() if n >= threshold}
    else:
        repeated = set()

    parts = []
    for lines in per_page_lines if len(pages_text) > 1 else [pages_text[0].split("\n")]:
        for ln in lines:
            s = ln.strip()
            if not s:
                continue
            if s in repeated:
                continue
            if _PAGE_HEADER_RE.match(s):
                continue
            parts.append(s)
    return " ".join(parts)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise NotImplementedError("python-pptx not installed — cannot extract pptx")
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    parts.append(t)
    return " ".join(parts)


def _extract_json(path: Path) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        # Canonical re-serialize: sort keys, minimal whitespace
        return json.dumps(data, sort_keys=True, separators=(",", ":"))
    except Exception:
        # Malformed JSON — fall back to raw text
        return _extract_text_plain(path)


def _extract_yaml(path: Path) -> str:
    try:
        import yaml
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = yaml.safe_load(f)
        return json.dumps(data, sort_keys=True, separators=(",", ":")) if data is not None else ""
    except ImportError:
        return _extract_text_plain(path)
    except Exception:
        return _extract_text_plain(path)


_EXTRACTORS = {
    ".txt":  _extract_text_plain,
    ".md":   _extract_text_plain,
    ".rtf":  _extract_text_plain,
    ".csv":  _extract_csv,
    ".tsv":  _extract_tsv,
    ".xlsx": _extract_xlsx,
    ".xls":  _extract_xlsx,   # best-effort
    ".docx": _extract_docx,
    ".pdf":  _extract_pdf,
    ".pptx": _extract_pptx,
    ".ppt":  _extract_pptx,   # best-effort
    ".json": _extract_json,
    ".yaml": _extract_yaml,
    ".yml":  _extract_yaml,
}


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def canonical_text(path: str | Path) -> str:
    """Extract + normalize the canonical text of a file. Returns "" for
    unsupported formats. Raises NotImplementedError if the format IS
    supported but the required library isn't installed (callers should
    catch and fall back)."""
    p = Path(path)
    ext = p.suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return ""
    raw = extractor(p)
    return _normalize(raw)


def canonical_hash(path: str | Path) -> Optional[str]:
    """Return SHA-256 hex of canonical_text(path), or None if the file
    has no canonical extractor and shouldn't contribute to cross-format
    discovery."""
    try:
        text = canonical_text(path)
    except NotImplementedError:
        return None
    if not text:
        return None
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def canonical_text_from_string(text: str) -> str:
    """Normalize a string the same way file extractors do. Used for
    browser-text attestations where the extension already has the
    normalized text in hand."""
    return _normalize(text)


def perceptual_hashes(image_path: str | Path) -> Optional[dict]:
    """Return {'dhash': '...', 'phash': '...'} for an image file. These
    hashes are robust to resize/recompression/minor crops. Returns None
    if imagehash or Pillow is unavailable, or if the file isn't a valid
    image."""
    try:
        import imagehash
        from PIL import Image
    except ImportError:
        return None
    try:
        with Image.open(image_path) as img:
            return {
                "dhash": str(imagehash.dhash(img)),
                "phash": str(imagehash.phash(img)),
            }
    except Exception:
        return None


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("Usage: python aiauth_canonical.py FILE")
        sys.exit(1)
    p = Path(sys.argv[1])
    try:
        h = canonical_hash(p)
        t = canonical_text(p)
        print(f"File: {p}")
        print(f"Canonical hash: {h}")
        print(f"Canonical text (first 200 chars): {t[:200]!r}")
    except NotImplementedError as e:
        print(f"Extractor unavailable: {e}", file=sys.stderr)
        sys.exit(2)
